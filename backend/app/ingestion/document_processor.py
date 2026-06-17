import io
import os
import uuid
from datetime import datetime

from langchain_chroma import Chroma
from langchain_core.documents import Document
from langchain_openai import OpenAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader

CHROMA_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "chroma_db")
COLLECTION_NAME = "documents"

_vector_store: Chroma | None = None


def _get_embeddings() -> OpenAIEmbeddings:
    return OpenAIEmbeddings(model="text-embedding-3-small")


def _get_store() -> Chroma:
    global _vector_store
    if _vector_store is None:
        _vector_store = Chroma(
            collection_name=COLLECTION_NAME,
            embedding_function=_get_embeddings(),
            persist_directory=os.path.abspath(CHROMA_DIR),
        )
    return _vector_store


def extract_text_from_pdf(file_bytes: bytes) -> list[dict]:
    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            pages.append({"page": i + 1, "text": text.strip()})
    return pages


def extract_text_from_pptx(file_bytes: bytes) -> list[dict]:

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        text = "\n".join(parts)
        if text.strip():
            slides.append({"page": i + 1, "text": text.strip()})
    return slides


def process_and_store_document(filename: str, file_bytes: bytes) -> dict:
    doc_id = str(uuid.uuid4())
    ext = os.path.splitext(filename)[1].lower()

    if ext == ".pdf":
        pages = extract_text_from_pdf(file_bytes)
    elif ext in (".pptx", ".ppt"):
        pages = extract_text_from_pptx(file_bytes)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    if not pages:
        return {"doc_id": doc_id, "filename": filename, "chunks": 0}

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    documents = []
    for page in pages:
        chunks = splitter.split_text(page["text"])
        for idx, chunk in enumerate(chunks):
            documents.append(
                Document(
                    page_content=chunk,
                    metadata={
                        "doc_id": doc_id,
                        "filename": filename,
                        "page_or_slide": page["page"],
                        "chunk_index": idx,
                        "uploaded_at": datetime.utcnow().isoformat(),
                    },
                )
            )

    store = _get_store()
    store.add_documents(documents)

    return {"doc_id": doc_id, "filename": filename, "chunks": len(documents)}


def search_documents(query: str, k: int = 5) -> list[Document]:
    store = _get_store()
    return store.similarity_search(query, k=k)


def list_documents() -> list[dict]:
    store = _get_store()
    data = store.get(include=["metadatas"])
    seen = {}
    for meta in data["metadatas"]:
        did = meta.get("doc_id")
        if not did:
            continue
        if did not in seen:
            seen[did] = {
                "doc_id": did,
                "filename": meta.get("filename", "unknown"),
                "uploaded_at": meta.get("uploaded_at", ""),
                "chunks": 0,
            }
        seen[did]["chunks"] += 1
    return list(seen.values())


def delete_document(doc_id: str) -> bool:
    store = _get_store()
    data = store.get()
    ids_to_delete = []
    for i, meta in enumerate(data["metadatas"]):
        if meta.get("doc_id") == doc_id:
            ids_to_delete.append(data["ids"][i])
    if ids_to_delete:
        store._collection.delete(ids_to_delete)
        return True
    return False
